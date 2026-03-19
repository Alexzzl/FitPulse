from __future__ import annotations

import os
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
APP_NAME = os.environ.get("FITPULSE_APP_NAME", "FitPulse TV")
COMPANY_NAME = os.environ.get("FITPULSE_COMPANY_NAME", "Your Company Name")
SUPPORT_EMAIL = os.environ.get("FITPULSE_SUPPORT_EMAIL", "your-support@example.com")
SUPPORT_URL = os.environ.get("FITPULSE_SUPPORT_URL", "https://your-domain.example")
SELLER_ADDRESS = os.environ.get("FITPULSE_SELLER_ADDRESS", "Your Address")
OUTPUT_DIR = ROOT / "docs" / "outputs"
OUTPUT = OUTPUT_DIR / "App_Description_FitPulse_TV_filled.pptx"
SCREENSHOT_DIR = ROOT / "store-screenshots"
UI_DIR = ROOT / "ui-description-screenshots"


def resolve_template() -> Path:
    candidates = [
        os.environ.get("FITPULSE_PPT_TEMPLATE"),
        ROOT / "App_Description_template_eng_v.1.42.pptx",
        ROOT.parent / "MoneTV" / "App_Description_template_eng_v.1.42.pptx",
    ]

    for candidate in candidates:
        if not candidate:
            continue
        path = Path(candidate)
        if path.exists():
            return path

    raise FileNotFoundError(
        "PPT template not found. Set FITPULSE_PPT_TEMPLATE or place "
        "App_Description_template_eng_v.1.42.pptx in FitPulse0 or the sibling MoneTV project."
    )


def require_file(path: Path) -> Path:
    if not path.exists():
        raise FileNotFoundError(f"Required asset not found: {path}")
    return path


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear_and_set_text(shape, lines: Iterable[str], font_size: int = 20, bold: bool = False) -> None:
    text_frame = shape.text_frame
    text_frame.clear()

    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def set_cell_text(cell, text: str, font_size: int = 12, bold: bool = False) -> None:
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def find_font(size: int):
    candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def create_collage(items: list[tuple[Path, str]], output_path: Path, title: str) -> None:
    width, height = 1600, 900
    margin = 40
    gutter = 24
    title_height = 70
    cols, rows = 2, 2
    cell_width = (width - margin * 2 - gutter) // cols
    cell_height = (height - margin * 2 - title_height - gutter) // rows
    image_height = cell_height - 52

    canvas = Image.new("RGB", (width, height), (248, 248, 248))
    draw = ImageDraw.Draw(canvas)
    title_font = find_font(34)
    label_font = find_font(24)

    draw.text((margin, 18), title, fill=(28, 28, 28), font=title_font)

    for idx, (image_path, label) in enumerate(items):
        row = idx // cols
        col = idx % cols
        x = margin + col * (cell_width + gutter)
        y = margin + title_height + row * (cell_height + gutter)

        image = Image.open(image_path).convert("RGB")
        fitted = ImageOps.fit(image, (cell_width, image_height), method=Image.Resampling.LANCZOS)
        canvas.paste(fitted, (x, y))

        label_top = y + image_height
        draw.rectangle((x, label_top, x + cell_width, y + cell_height), fill=(255, 255, 255))
        draw.rectangle((x, y, x + cell_width, y + cell_height), outline=(210, 210, 210), width=2)
        draw.text((x + 16, label_top + 12), label, fill=(40, 40, 40), font=label_font)

    canvas.save(output_path, quality=90)


def add_centered_box(slide, left, top, width, height, text: str, fill_rgb: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_arrow_text(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def build_collages() -> tuple[Path, Path]:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    UI_DIR.mkdir(parents=True, exist_ok=True)

    usage_collage = UI_DIR / "10-usage-scenario-collage.jpg"
    menu_collage = UI_DIR / "11-menu-function-collage.jpg"

    create_collage(
        [
            (require_file(SCREENSHOT_DIR / "01-home-dashboard.jpg"), "1. Home dashboard"),
            (require_file(UI_DIR / "08-day-detail.jpg"), "2. Day detail"),
            (require_file(SCREENSHOT_DIR / "04-workout-player.jpg"), "3. Workout player"),
            (require_file(UI_DIR / "09-workout-complete.jpg"), "4. Workout complete"),
        ],
        usage_collage,
        f"{APP_NAME} - Usage Scenario",
    )

    create_collage(
        [
            (require_file(SCREENSHOT_DIR / "02-library-explore.jpg"), "Library"),
            (require_file(SCREENSHOT_DIR / "03-classic-programs.jpg"), "Classic Programs"),
            (require_file(UI_DIR / "07-plan-calendar.jpg"), "Plan Calendar"),
            (require_file(UI_DIR / "06-history-screen.jpg"), "History"),
        ],
        menu_collage,
        f"{APP_NAME} - Main Screens and Functions",
    )

    return usage_collage, menu_collage


def main() -> None:
    template = resolve_template()
    usage_collage, menu_collage = build_collages()

    prs = Presentation(str(template))
    remove_slide(prs, 0)

    cover = prs.slides[0]
    cover.shapes[0].text = f"{APP_NAME} Description"
    cover.shapes[1].text = COMPANY_NAME

    revision = prs.slides[1]
    revision.shapes[0].text = "Revision History"
    table = revision.shapes[1].table
    set_cell_text(table.cell(0, 0), "Version", 12, True)
    set_cell_text(table.cell(0, 1), "Date", 12, True)
    set_cell_text(table.cell(0, 2), "Description", 12, True)
    set_cell_text(table.cell(0, 3), "Author", 12, True)
    set_cell_text(table.cell(1, 0), "1.0")
    set_cell_text(table.cell(1, 1), "March 19, 2026")
    set_cell_text(table.cell(1, 2), f"Initial Samsung TV Seller Office UI description for {APP_NAME}.")
    set_cell_text(table.cell(1, 3), COMPANY_NAME)
    for col in range(4):
        set_cell_text(table.cell(2, col), "")

    contents = prs.slides[2]
    contents.shapes[0].text = "Contents"
    clear_and_set_text(
        contents.shapes[1],
        [
            "UI Structure",
            "Usage Scenario",
            "Menu & Function Description",
            "Key Policy",
            "Additional Notes",
        ],
        font_size=24,
    )

    ui_overview = prs.slides[3]
    ui_overview.shapes[0].text = "UI Structure"
    clear_and_set_text(
        ui_overview.shapes[1],
        [
            "Whole UI Structure",
            "",
            f"{APP_NAME} is a Samsung TV fitness application for browsing workouts,",
            "opening guided training plans, and starting remote-friendly workout sessions.",
            "Only currently available user flows are described in this document.",
            "",
            "Primary user flow:",
            "Welcome -> Profile Ready -> Home",
            "Home -> Classic -> Plan Calendar -> Day Detail -> Workout Player",
            "Home / Me -> History",
            "",
            "Main screens in the current release:",
            "- Home dashboard",
            "- Library by body part",
            "- Classic programs",
            "- Plan calendar and day detail",
            "- Workout player and completion summary",
        ],
        font_size=18,
    )

    flow_slide = prs.slides[4]
    flow_slide.shapes[0].text = "UI Structure - Flow Graph"
    remove_shape(flow_slide.shapes[1])
    add_centered_box(flow_slide, Inches(1.0), Inches(1.4), Inches(1.6), Inches(0.65), "Welcome", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(2.75), Inches(1.45), Inches(0.45), Inches(0.4), "->")
    add_centered_box(flow_slide, Inches(3.2), Inches(1.4), Inches(1.85), Inches(0.65), "Home", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(5.15), Inches(1.45), Inches(0.45), Inches(0.4), "->")
    add_centered_box(flow_slide, Inches(5.65), Inches(1.4), Inches(1.8), Inches(0.65), "Library", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(3.95), Inches(2.2), Inches(0.3), Inches(0.4), "v")
    add_centered_box(flow_slide, Inches(3.0), Inches(2.65), Inches(2.3), Inches(0.68), "Classic Programs", (0xE0, 0x8A, 0x3B))
    add_arrow_text(flow_slide, Inches(4.0), Inches(3.45), Inches(0.3), Inches(0.4), "v")
    add_centered_box(flow_slide, Inches(2.85), Inches(3.9), Inches(2.6), Inches(0.68), "Plan Calendar", (0xE0, 0x8A, 0x3B))
    add_arrow_text(flow_slide, Inches(4.0), Inches(4.7), Inches(0.3), Inches(0.4), "v")
    add_centered_box(flow_slide, Inches(3.0), Inches(5.1), Inches(2.3), Inches(0.68), "Day Detail", (0xB9, 0x4E, 0x5D))
    add_arrow_text(flow_slide, Inches(5.55), Inches(5.2), Inches(0.45), Inches(0.4), "->")
    add_centered_box(flow_slide, Inches(6.05), Inches(5.1), Inches(2.0), Inches(0.68), "Workout Player", (0xB9, 0x4E, 0x5D))

    depth_slide = prs.slides[5]
    depth_slide.shapes[0].text = "UI Structure - Depth Navigation"
    remove_shape(depth_slide.shapes[1])
    textbox = depth_slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(7.3), Inches(4.5))
    clear_and_set_text(
        textbox,
        [
            "Depth 1 : Welcome, Profile Ready, Home, Library, Classic, Me",
            "Depth 2 : History, Plan Calendar, Day Detail",
            "Depth 3 : Get Ready, Workout Player, Rest, Workout Complete",
            "",
            "Global UI elements:",
            "- Sidebar navigation on main browsing screens",
            "- Directional focus highlight",
            "",
            "Contextual UI elements:",
            "- Workout plan cards",
            "- Day grid in the plan calendar",
            "- Start button on day detail page",
            "- Timer and action button in workout flow",
        ],
        font_size=20,
    )

    usage_summary = prs.slides[6]
    usage_summary.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_summary.shapes[1],
        [
            "Use Case 1: Browse a workout plan and start a guided session.",
            "",
            "1. Launch the app and go to the Home screen.",
            "2. Open Library or Classic Programs with the remote.",
            "3. Select a workout plan from Classic Programs.",
            "4. Choose an available day from the plan calendar.",
            "5. Review the workout summary on the Day Detail screen.",
            "6. Start the workout and continue through the workout player flow.",
            "7. Return to previous screens with the Back key.",
            "",
            "No account sign-in, activation, or purchase is required in the current release.",
        ],
        font_size=18,
    )

    usage_sample = prs.slides[7]
    usage_sample.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_sample.shapes[1],
        [
            "Use Case Title: Browse and start a workout session",
            "Path: Home -> Classic -> Plan Calendar -> Day Detail -> Workout Player",
            "",
            "Supporting notes:",
            "- Users can also browse workouts from Library.",
            "- The TV remote is used for directional focus navigation.",
            "- The completion summary appears after the guided workout flow.",
            "",
            "No account login is required.",
            "No in-app purchase is required.",
        ],
        font_size=18,
    )
    usage_picture = usage_sample.shapes[2]
    left, top, width, height = usage_picture.left, usage_picture.top, usage_picture.width, usage_picture.height
    remove_shape(usage_picture)
    usage_sample.shapes.add_picture(str(usage_collage), left, top, width=width, height=height)

    menu_slide = prs.slides[8]
    menu_slide.shapes[0].text = "Menu & Function Description"
    clear_and_set_text(
        menu_slide.shapes[1],
        [
            "Screen descriptions in the current release:",
            "",
            "Home: weekly activity, resume card, and recommended workouts.",
            "Library: body-part filter chips and workout cards.",
            "Classic Programs: structured plan cards for long-form training.",
            "Plan Calendar: day selection for a multi-day workout plan.",
            "History: previous completed sessions.",
            "",
            "Only currently implemented screens and functions are described in this document.",
        ],
        font_size=18,
    )
    menu_picture = menu_slide.shapes[2]
    left, top, width, height = menu_picture.left, menu_picture.top, menu_picture.width, menu_picture.height
    remove_shape(menu_picture)
    menu_slide.shapes.add_picture(str(menu_collage), left, top, width=width, height=height)

    key_slide = prs.slides[9]
    key_slide.shapes[0].text = "Key Policy"
    clear_and_set_text(
        key_slide.shapes[1],
        [
            "Standard Samsung TV navigation behavior is used in the current release.",
            "The app focuses on directional focus movement and Back key behavior.",
        ],
        font_size=18,
    )
    key_table = key_slide.shapes[2].table
    rows = [
        ("Button", "Action", "Remarks"),
        ("ENTER", "Select the focused card, button, or action", "Standard selection behavior"),
        ("UP / DOWN", "Move focus vertically", "Stops at the screen edge"),
        ("LEFT / RIGHT", "Move focus horizontally", "Stops at the screen edge"),
        ("BACK / RETURN", "Go to the previous screen", "Standard Samsung TV behavior"),
        ("PLAY / PAUSE", "Control timer-based workout flow when applicable", "Current workout flow only"),
        ("EXIT", "Close the application", "Standard Samsung TV behavior"),
        ("COLOR KEYS", "N/R", "No custom key mapping"),
        ("CHANNEL / NUMBER KEYS", "N/R", "No custom key mapping"),
        ("FAST FORWARD / REWIND", "N/R", "No custom key mapping"),
    ]
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            set_cell_text(key_table.cell(r_idx, c_idx), value, 11, r_idx == 0)

    notes_slide = prs.slides[10]
    notes_slide.shapes[0].text = "Additional Notes"
    remove_shape(notes_slide.shapes[1])
    textbox = notes_slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(7.5), Inches(4.6))
    clear_and_set_text(
        textbox,
        [
            "Language support",
            "This version supports English only. No in-app language change menu is provided.",
            "",
            "Seller / Support",
            f"Seller: {COMPANY_NAME}",
            f"Support E-mail: {SUPPORT_EMAIL}",
            f"Support URL: {SUPPORT_URL}",
            f"Seller Address: {SELLER_ADDRESS}",
            "",
            "Additional notes",
            "This app is a Samsung TV web application package (.wgt).",
            "The current release focuses on workout browsing, plan selection, and guided remote-based workout flow.",
        ],
        font_size=18,
    )

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"template={template}")
    print(f"saved={OUTPUT}")


if __name__ == "__main__":
    main()
